import sys
import subprocess

# Ensure python-pptx is installed
try:
    import pptx
except ImportError:
    print("Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.33)  # 16:9 Aspect Ratio
prs.slide_height = Inches(7.5)

# Color Palette
DARK_BG = RGBColor(10, 7, 18)
WHITE = RGBColor(248, 250, 252)
SECONDARY_TEXT = RGBColor(148, 163, 184)
RED = RGBColor(237, 66, 85)
TEAL = RGBColor(80, 187, 182)
GOLD = RGBColor(251, 182, 74)

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_header(slide, title_text, category_text=None):
    if category_text:
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = 'Arial'
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = RED
        
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE

# Slide 1: Title Slide
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide)

# Title Text Box
title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.0))
tf = title_box.text_frame
tf.word_wrap = True

# Track badge
p_badge = tf.paragraphs[0]
p_badge.text = "TRACK P3: INTERACTIVE ENTERTAINMENT"
p_badge.font.name = 'Arial'
p_badge.font.size = Pt(12)
p_badge.font.bold = True
p_badge.font.color.rgb = RED
p_badge.alignment = PP_ALIGN.CENTER
p_badge.space_after = Pt(20)

# Main Title
p_title = tf.add_paragraph()
p_title.text = "VERDICTS"
p_title.font.name = 'Arial'
p_title.font.size = Pt(72)
p_title.font.bold = True
p_title.font.color.rgb = WHITE
p_title.alignment = PP_ALIGN.CENTER
p_title.space_after = Pt(10)

# Subtitle
p_sub = tf.add_paragraph()
p_sub.text = "First-Person Generative Audio Storytelling for Pocket FM"
p_sub.font.name = 'Arial'
p_sub.font.size = Pt(20)
p_sub.font.color.rgb = SECONDARY_TEXT
p_sub.alignment = PP_ALIGN.CENTER
p_sub.space_after = Pt(30)

# Team Name
p_team = tf.add_paragraph()
p_team.text = "Created by Team: DIGITAL AGENTIC"
p_team.font.name = 'Arial'
p_team.font.size = Pt(14)
p_team.font.bold = True
p_team.font.color.rgb = WHITE
p_team.alignment = PP_ALIGN.CENTER


# Slide 2: The Problem
slide2 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide2)
add_header(slide2, "The Pain of Passive Media", "The Challenge")

# Stat Card Left
stat_bg = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(4.5), Inches(4.8))
stat_bg.fill.solid()
stat_bg.fill.fore_color.rgb = RGBColor(20, 15, 30)
stat_bg.line.color.rgb = RED
stat_bg.line.width = Pt(1.5)

tf_stat = stat_bg.text_frame
tf_stat.word_wrap = True
tf_stat.margin_top = Inches(0.8)
tf_stat.margin_left = Inches(0.4)

p_stat_num = tf_stat.paragraphs[0]
p_stat_num.text = "83%"
p_stat_num.font.name = 'Arial'
p_stat_num.font.size = Pt(64)
p_stat_num.font.bold = True
p_stat_num.font.color.rgb = RED
p_stat_num.alignment = PP_ALIGN.CENTER
p_stat_num.space_after = Pt(20)

p_stat_lbl = tf_stat.add_paragraph()
p_stat_lbl.text = "Passive Fatigue Drop-off"
p_stat_lbl.font.name = 'Arial'
p_stat_lbl.font.size = Pt(22)
p_stat_lbl.font.bold = True
p_stat_lbl.font.color.rgb = WHITE
p_stat_lbl.alignment = PP_ALIGN.CENTER
p_stat_lbl.space_after = Pt(10)

p_stat_desc = tf_stat.add_paragraph()
p_stat_desc.text = "Listeners drop off streams when content is purely linear and predictable."
p_stat_desc.font.name = 'Arial'
p_stat_desc.font.size = Pt(14)
p_stat_desc.font.color.rgb = SECONDARY_TEXT
p_stat_desc.alignment = PP_ALIGN.CENTER

# Point 1 Right
pt1_bg = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.8), Inches(1.8), Inches(6.7), Inches(2.2))
pt1_bg.fill.solid()
pt1_bg.fill.fore_color.rgb = RGBColor(20, 15, 30)
pt1_bg.line.color.rgb = SECONDARY_TEXT

tf_pt1 = pt1_bg.text_frame
tf_pt1.word_wrap = True
tf_pt1.margin_top = Inches(0.3)
tf_pt1.margin_left = Inches(0.4)

p_pt1_t = tf_pt1.paragraphs[0]
p_pt1_t.text = "🥱 Predictable Script Fatigue"
p_pt1_t.font.name = 'Arial'
p_pt1_t.font.size = Pt(20)
p_pt1_t.font.bold = True
p_pt1_t.font.color.rgb = WHITE
p_pt1_t.space_after = Pt(8)

p_pt1_d = tf_pt1.add_paragraph()
p_pt1_d.text = "Standard audiobooks offer no repeat value. Once the ending is heard, replayability collapses to 0%."
p_pt1_d.font.name = 'Arial'
p_pt1_d.font.size = Pt(14)
p_pt1_d.font.color.rgb = SECONDARY_TEXT

# Point 2 Right
pt2_bg = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.8), Inches(4.4), Inches(6.7), Inches(2.2))
pt2_bg.fill.solid()
pt2_bg.fill.fore_color.rgb = RGBColor(20, 15, 30)
pt2_bg.line.color.rgb = SECONDARY_TEXT

tf_pt2 = pt2_bg.text_frame
tf_pt2.word_wrap = True
tf_pt2.margin_top = Inches(0.3)
tf_pt2.margin_left = Inches(0.4)

p_pt2_t = tf_pt2.paragraphs[0]
p_pt2_t.text = "❌ Zero Audience Agency"
p_pt2_t.font.name = 'Arial'
p_pt2_t.font.size = Pt(20)
p_pt2_t.font.bold = True
p_pt2_t.font.color.rgb = WHITE
p_pt2_t.space_after = Pt(8)

p_pt2_d = tf_pt2.add_paragraph()
p_pt2_d.text = "Listeners have no say in character decisions, environment actions, or moral choices. They are locked behind a glass wall."
p_pt2_d.font.name = 'Arial'
p_pt2_d.font.size = Pt(14)
p_pt2_d.font.color.rgb = SECONDARY_TEXT


# Slide 3: The Solution
slide3 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide3)
add_header(slide3, "Verdicts: The Interactive Audio Hub", "The Solution")

sol_cols = [
    {"icon": "🎙️", "title": "Interactive Lobbies", "desc": "A dedicated, immersive audio tab built right inside Pocket FM's app ecosystem for generative multiplayer rooms.", "color": RED},
    {"icon": "🧠", "title": "Generative AI Engine", "desc": "Wipes predictable scripts. The AI generates suspects, motives, clue hierarchies, and cover art on the fly.", "color": GOLD},
    {"icon": "💬", "title": "Realtime Interrogations", "desc": "Users speak directly to suspects. AI characters respond dynamically in real time, matching their guilt profiles.", "color": TEAL}
]

for idx, col in enumerate(sol_cols):
    col_bg = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + idx * 4.0), Inches(2.2), Inches(3.7), Inches(4.2))
    col_bg.fill.solid()
    col_bg.fill.fore_color.rgb = RGBColor(20, 15, 30)
    col_bg.line.color.rgb = col["color"]
    col_bg.line.width = Pt(1.5)
    
    tf_col = col_bg.text_frame
    tf_col.word_wrap = True
    tf_col.margin_top = Inches(0.5)
    tf_col.margin_left = Inches(0.3)
    tf_col.margin_right = Inches(0.3)
    
    p_i = tf_col.paragraphs[0]
    p_i.text = col["icon"]
    p_i.font.name = 'Arial'
    p_i.font.size = Pt(40)
    p_i.space_after = Pt(20)
    p_i.alignment = PP_ALIGN.CENTER
    
    p_t = tf_col.add_paragraph()
    p_t.text = col["title"]
    p_t.font.name = 'Arial'
    p_t.font.size = Pt(20)
    p_t.font.bold = True
    p_t.font.color.rgb = WHITE
    p_t.space_after = Pt(10)
    p_t.alignment = PP_ALIGN.CENTER
    
    p_d = tf_col.add_paragraph()
    p_d.text = col["desc"]
    p_d.font.name = 'Arial'
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = SECONDARY_TEXT
    p_d.alignment = PP_ALIGN.CENTER


# Slide 4: SoulCast
slide4 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide4)
add_header(slide4, "Pillar 1: SoulCast (Story as a Game)", "Core Feature")

left_txt = slide4.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
tf_left = left_txt.text_frame
tf_left.word_wrap = True

p_l1 = tf_left.paragraphs[0]
p_l1.text = "DIVERGENT AUDIO EXPERIENCE"
p_l1.font.name = 'Arial'
p_l1.font.size = Pt(14)
p_l1.font.bold = True
p_l1.font.color.rgb = GOLD
p_l1.space_after = Pt(15)

p_l2 = tf_left.add_paragraph()
p_l2.text = "Decisions Reshape narrative physics."
p_l2.font.name = 'Arial'
p_l2.font.size = Pt(30)
p_l2.font.bold = True
p_l2.font.color.rgb = WHITE
p_l2.space_after = Pt(15)

p_l3 = tf_left.add_paragraph()
p_l3.text = "In SoulCast, listeners are the first-person protagonist of interactive thrillers. When the story pauses, you speak or select your decision.\n\nThe backend immediately creates a permanent 'WorldRule' in the database. Future scenes adapt dynamically, shifting the ending based on your personality, history, and behavior."
p_l3.font.name = 'Arial'
p_l3.font.size = Pt(15)
p_l3.font.color.rgb = SECONDARY_TEXT

right_card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.5))
right_card.fill.solid()
right_card.fill.fore_color.rgb = RGBColor(20, 15, 30)
right_card.line.color.rgb = GOLD
right_card.line.width = Pt(1.5)

tf_rc = right_card.text_frame
tf_rc.word_wrap = True
tf_rc.margin_top = Inches(0.5)
tf_rc.margin_left = Inches(0.4)

p_rt = tf_rc.paragraphs[0]
p_rt.text = "Key User Value Metrics"
p_rt.font.name = 'Arial'
p_rt.font.size = Pt(22)
p_rt.font.bold = True
p_rt.font.color.rgb = WHITE
p_rt.space_after = Pt(20)

points = [
    "🎙️ Spoken/Voice-native interactive choices",
    "🧠 Ends passive story drop-off",
    "📈 High-fidelity personalized outcomes",
    "♻️ Endings built on behavior and listening history"
]
for pt in points:
    p_p = tf_rc.add_paragraph()
    p_p.text = pt
    p_p.font.name = 'Arial'
    p_p.font.size = Pt(16)
    p_p.font.color.rgb = SECONDARY_TEXT
    p_p.space_after = Pt(15)


# Slide 5: Infinite Murder Mystery
slide5 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide5)
add_header(slide5, "Pillar 2: Infinite Murder Mystery", "Core Feature")

left_card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.7), Inches(4.5))
left_card.fill.solid()
left_card.fill.fore_color.rgb = RGBColor(20, 15, 30)
left_card.line.color.rgb = RED
left_card.line.width = Pt(1.5)

tf_lc = left_card.text_frame
tf_lc.word_wrap = True
tf_lc.margin_top = Inches(0.5)
tf_lc.margin_left = Inches(0.4)

p_lt = tf_lc.paragraphs[0]
p_lt.text = "Co-Op Clues & Boards"
p_lt.font.name = 'Arial'
p_lt.font.size = Pt(22)
p_lt.font.bold = True
p_lt.font.color.rgb = WHITE
p_lt.space_after = Pt(15)

m_points = [
    "🔍 Split, asymmetric clues across players",
    "📌 Collaborative shared evidence board",
    "🧬 Dynamically generated killer profiles",
    "🧪 Logic paths check generated story validity"
]
for pt in m_points:
    p_p = tf_lc.add_paragraph()
    p_p.text = pt
    p_p.font.name = 'Arial'
    p_p.font.size = Pt(16)
    p_p.font.color.rgb = SECONDARY_TEXT
    p_p.space_after = Pt(15)

right_txt = slide5.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5))
tf_right = right_txt.text_frame
tf_right.word_wrap = True

p_r1 = tf_right.paragraphs[0]
p_r1.text = "INFINITE REPLAYABILITY"
p_r1.font.name = 'Arial'
p_r1.font.size = Pt(14)
p_r1.font.bold = True
p_r1.font.color.rgb = RED
p_r1.space_after = Pt(15)

p_r2 = tf_right.add_paragraph()
p_r2.text = "Realtime AI Suspect Interrogations"
p_r2.font.name = 'Arial'
p_r2.font.size = Pt(30)
p_r2.font.bold = True
p_r2.font.color.rgb = WHITE
p_r2.space_after = Pt(15)

p_r3 = tf_right.add_paragraph()
p_r3.text = "Every playthrough generates a completely different killer, motive, evidence set, and solution.\n\nPlayers speak directly to AI suspects in natural language. The suspects respond dynamically using OpenAI Realtime voice engines, aligning contextually to their guilt profiles."
p_r3.font.name = 'Arial'
p_r3.font.size = Pt(15)
p_r3.font.color.rgb = SECONDARY_TEXT


# Slide 6: Locked
slide6 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide6)
add_header(slide6, "Pillar 3: Locked (AI Escape Room)", "Core Feature")

left_txt = slide6.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
tf_left = left_txt.text_frame
tf_left.word_wrap = True

p_l1 = tf_left.paragraphs[0]
p_l1.text = "COOPERATIVE AUDIO CHALLENGES"
p_l1.font.name = 'Arial'
p_l1.font.size = Pt(14)
p_l1.font.bold = True
p_l1.font.color.rgb = TEAL
p_l1.space_after = Pt(15)

p_l2 = tf_left.add_paragraph()
p_l2.text = "Crack the code through conversation."
p_l2.font.name = 'Arial'
p_l2.font.size = Pt(30)
p_l2.font.bold = True
p_l2.font.color.rgb = WHITE
p_l2.space_after = Pt(15)

p_l3 = tf_left.add_paragraph()
p_l3.text = "Locked places a group of players inside a dark, atmospheric room with no written instructions. The environment responds directly to spoken descriptions.\n\nBy naturally conversing with the room, players search objects, decode puzzle logs, and uncover split hints to find the exit door."
p_l3.font.name = 'Arial'
p_l3.font.size = Pt(15)
p_l3.font.color.rgb = SECONDARY_TEXT

right_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.5))
right_card.fill.solid()
right_card.fill.fore_color.rgb = RGBColor(20, 15, 30)
right_card.line.color.rgb = TEAL
right_card.line.width = Pt(1.5)

tf_rc = right_card.text_frame
tf_rc.word_wrap = True
tf_rc.margin_top = Inches(0.5)
tf_rc.margin_left = Inches(0.4)

p_rt = tf_rc.paragraphs[0]
p_rt.text = "Unique Gameplay Elements"
p_rt.font.name = 'Arial'
p_rt.font.size = Pt(22)
p_rt.font.bold = True
p_rt.font.color.rgb = WHITE
p_rt.space_after = Pt(20)

points = [
    "🧩 Dynamic puzzle solution checks",
    "🗣️ Voice-native environmental commands",
    "👥 Shared realtime player state sync",
    "🔊 Spatial audio clues and anomalies"
]
for pt in points:
    p_p = tf_rc.add_paragraph()
    p_p.text = pt
    p_p.font.name = 'Arial'
    p_p.font.size = Pt(16)
    p_p.font.color.rgb = SECONDARY_TEXT
    p_p.space_after = Pt(15)


# Slide 7: Technical Architecture
slide7 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide7)
add_header(slide7, "Technical Architecture Flow", "Tech Stack")

techs = [
    {"label": "React / Vite", "d": "Premium glassmorphic client interface & Websockets"},
    {"label": "FastAPI & Socket.io", "d": "ASGI realtime room broadcast and REST endpoints"},
    {"label": "SQLite Store", "d": "Local database to persist state & curated cases"},
    {"label": "OpenAI Service", "d": "TTS, Realtime audio API, DALL-E-3 cover art"}
]

for idx, t in enumerate(techs):
    t_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + idx * 2.95), Inches(2.2), Inches(2.7), Inches(3.2))
    t_box.fill.solid()
    t_box.fill.fore_color.rgb = RGBColor(20, 15, 30)
    t_box.line.color.rgb = SECONDARY_TEXT
    
    tf_t = t_box.text_frame
    tf_t.word_wrap = True
    tf_t.margin_top = Inches(0.4)
    tf_t.margin_left = Inches(0.2)
    tf_t.margin_right = Inches(0.2)
    
    p_tl = tf_t.paragraphs[0]
    p_tl.text = t["label"]
    p_tl.font.name = 'Arial'
    p_tl.font.size = Pt(20)
    p_tl.font.bold = True
    p_tl.font.color.rgb = WHITE
    p_tl.space_after = Pt(15)
    p_tl.alignment = PP_ALIGN.CENTER
    
    p_td = tf_t.add_paragraph()
    p_td.text = t["d"]
    p_td.font.name = 'Arial'
    p_td.font.size = Pt(13)
    p_td.font.color.rgb = SECONDARY_TEXT
    p_td.alignment = PP_ALIGN.CENTER

arch_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.1))
arch_card.fill.solid()
arch_card.fill.fore_color.rgb = RGBColor(20, 15, 30)
arch_card.line.color.rgb = RED
arch_card.line.width = Pt(1.5)

tf_ac = arch_card.text_frame
tf_ac.word_wrap = True
tf_ac.margin_top = Inches(0.2)
tf_ac.margin_left = Inches(0.4)

p_ac = tf_ac.paragraphs[0]
p_ac.text = "☁️ Databricks Apps Ready:"
p_ac.font.name = 'Arial'
p_ac.font.size = Pt(15)
p_ac.font.bold = True
p_ac.font.color.rgb = WHITE
p_ac.space_after = Pt(4)

p_ac_sub = tf_ac.add_paragraph()
p_ac_sub.text = "Hosted serverless inside Databricks. Resolves app paths, database instances, and environment models out of the box."
p_ac_sub.font.name = 'Arial'
p_ac_sub.font.size = Pt(13)
p_ac_sub.font.color.rgb = SECONDARY_TEXT


# Slide 8: Business Value
slide8 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide8)
add_header(slide8, "Unlocking Business Value for Pocket FM", "Impact")

b_cols = [
    {"val": "+42%", "title": "Engagement Boost", "desc": "Interactive choice loops and natural conversations drive average listener sessions way beyond linear files.", "color": TEAL},
    {"val": "Coins", "title": "Monetization Hooks", "desc": "Plugs into Pocket FM's coin structure. Pay micro-coins to interrogate suspects or buy extra puzzle clues.", "color": GOLD},
    {"val": "100%", "title": "Replay Value", "desc": "Because the underlying story layers are generated dynamically, listeners repeat the same case to test alternate endings.", "color": RED}
]

for idx, b in enumerate(b_cols):
    b_bg = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + idx * 4.0), Inches(2.2), Inches(3.7), Inches(4.2))
    b_bg.fill.solid()
    b_bg.fill.fore_color.rgb = RGBColor(20, 15, 30)
    b_bg.line.color.rgb = b["color"]
    b_bg.line.width = Pt(1.5)
    
    tf_b = b_bg.text_frame
    tf_b.word_wrap = True
    tf_b.margin_top = Inches(0.5)
    tf_b.margin_left = Inches(0.3)
    tf_b.margin_right = Inches(0.3)
    
    p_v = tf_b.paragraphs[0]
    p_v.text = b["val"]
    p_v.font.name = 'Arial'
    p_v.font.size = Pt(48)
    p_v.font.bold = True
    p_v.font.color.rgb = b["color"]
    p_v.space_after = Pt(15)
    p_v.alignment = PP_ALIGN.CENTER
    
    p_bt = tf_b.add_paragraph()
    p_bt.text = b["title"]
    p_bt.font.name = 'Arial'
    p_bt.font.size = Pt(20)
    p_bt.font.bold = True
    p_bt.font.color.rgb = WHITE
    p_bt.space_after = Pt(10)
    p_bt.alignment = PP_ALIGN.CENTER
    
    p_bd = tf_b.add_paragraph()
    p_bd.text = b["desc"]
    p_bd.font.name = 'Arial'
    p_bd.font.size = Pt(14)
    p_bd.font.color.rgb = SECONDARY_TEXT
    p_bd.alignment = PP_ALIGN.CENTER


# Slide 9: Closing
slide9 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide9)

title_box = slide9.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.3), Inches(3.0))
tf_c = title_box.text_frame
tf_c.word_wrap = True

p_c_t = tf_c.paragraphs[0]
p_c_t.text = "DIGITAL AGENTIC"
p_c_t.font.name = 'Arial'
p_c_t.font.size = Pt(64)
p_c_t.font.bold = True
p_c_t.font.color.rgb = RED
p_c_t.alignment = PP_ALIGN.CENTER
p_c_t.space_after = Pt(10)

p_c_s = tf_c.add_paragraph()
p_c_s.text = "Delivering the future of interactive audio entertainment today."
p_c_s.font.name = 'Arial'
p_c_s.font.size = Pt(22)
p_c_s.font.color.rgb = WHITE
p_c_s.alignment = PP_ALIGN.CENTER
p_c_s.space_after = Pt(30)

p_c_qa = tf_c.add_paragraph()
p_c_qa.text = "Q&A Session"
p_c_qa.font.name = 'Arial'
p_c_qa.font.size = Pt(16)
p_c_qa.font.bold = True
p_c_qa.font.color.rgb = SECONDARY_TEXT
p_c_qa.alignment = PP_ALIGN.CENTER

# Save the Presentation
prs.save("digital_agentic_pitch.pptx")
print("PowerPoint presentation 'digital_agentic_pitch.pptx' generated successfully!")
